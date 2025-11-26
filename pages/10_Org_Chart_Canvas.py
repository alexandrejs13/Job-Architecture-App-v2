import tempfile
from pathlib import Path
from typing import Dict, List

import pandas as pd
import streamlit as st
from pptx import Presentation
from pyvis.network import Network
import streamlit.components.v1 as components


@st.cache_data
def load_profiles(path: Path) -> pd.DataFrame:
    df = pd.read_excel(path).fillna("")
    needed = ["Job Family", "Sub Job Family", "Job Profile", "Global Grade", "Role Description"]
    missing = [c for c in needed if c not in df.columns]
    if missing:
        st.error(f"Colunas ausentes no Job Profile.xlsx: {', '.join(missing)}")
        return pd.DataFrame(columns=needed)
    return df


@st.cache_data
def list_ppts() -> List[Path]:
    candidates = []
    for folder in ["PPTs", "ppts"]:
        d = Path(folder)
        if d.exists():
            candidates.extend(d.glob("*.pptx"))
    return sorted(candidates)


def match_ppt(ppts: List[Path], family: str, sub_family: str, profile: str) -> List[Path]:
    targets = [profile.lower(), sub_family.lower(), family.lower()]
    hits = []
    for ppt in ppts:
        stem = ppt.stem.lower()
        if any(t and t in stem for t in targets):
            hits.append(ppt)
    return hits


def extract_first_slide(ppt: Path) -> Dict[str, str]:
    try:
        prs = Presentation(ppt)
        if not prs.slides:
            return {"title": ppt.stem, "body": ""}
        slide = prs.slides[0]
        title, body_parts = "", []
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            txt = shape.text.strip()
            if not txt:
                continue
            if not title:
                title = txt
            else:
                body_parts.append(txt)
        return {"title": title or ppt.stem, "body": "\n\n".join(body_parts)}
    except Exception:
        return {"title": ppt.stem, "body": ""}


def render_infocard(profile: Dict, attachments: List[Path]) -> None:
    st.markdown(f"### {profile['Job Profile']}")
    st.markdown(f"**Job Family:** {profile['Job Family']}  \n**Sub Job Family:** {profile['Sub Job Family']}  \n**Global Grade:** {profile.get('Global Grade', '—')}")
    desc = profile.get("Role Description", "")
    if desc:
        st.markdown(desc)
    if attachments:
        st.markdown("**Infográficos (PPT):**")
        cols = st.columns(len(attachments))
        for idx, ppt in enumerate(attachments):
            meta = extract_first_slide(ppt)
            with ppt.open("rb") as f:
                data = f.read()
            cols[idx].download_button(
                label=f"{meta['title'][:24]}... (baixar)",
                data=data,
                file_name=ppt.name,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                key=f"dl-{ppt.name}-{idx}",
            )
            if meta["body"]:
                cols[idx].markdown(f"<small>{meta['body'][:180]}...</small>", unsafe_allow_html=True)


def build_network(df: pd.DataFrame, selected_family: str | None = None) -> str:
    net = Network(height="600px", width="100%", directed=False, bgcolor="#f8f9fb", font_color="#111")
    net.toggle_physics(True)
    families = df["Job Family"].unique()
    for fam in families:
        if selected_family and fam != selected_family:
            continue
        net.add_node(f"fam:{fam}", label=fam, size=24, color="#145efc", shape="dot", title=f"Job Family: {fam}")
        sub_df = df[df["Job Family"] == fam]
        subs = sub_df["Sub Job Family"].unique()
        for sub in subs:
            net.add_node(f"sub:{fam}:{sub}", label=sub, size=18, color="#5a8bff", shape="dot", title=f"Sub Job Family: {sub}")
            net.add_edge(f"fam:{fam}", f"sub:{fam}:{sub}", color="#cbd5e1")
            profs = sub_df[sub_df["Sub Job Family"] == sub]["Job Profile"].unique()
            for prof in profs:
                net.add_node(
                    f"job:{fam}:{sub}:{prof}",
                    label=prof,
                    size=12,
                    color="#7c3aed",
                    shape="ellipse",
                    title=f"{prof} ({sub} · {fam})",
                )
                net.add_edge(f"sub:{fam}:{sub}", f"job:{fam}:{sub}:{prof}", color="#e2e8f0")

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".html")
    net.show(tmp.name)
    with open(tmp.name, "r", encoding="utf-8") as f:
        html = f.read()
    return html


def main():
    st.set_page_config(page_title="Org Chart Canvas", layout="wide")
    df = load_profiles(Path("data") / "Job Profile.xlsx")
    if df.empty:
        st.stop()
    ppts = list_ppts()

    st.title("Org Chart Canvas")
    st.markdown("Selecione uma Job Family e navegue no organograma; clique em um cargo na lista para ver o infográfico.")

    families = sorted(df["Job Family"].dropna().unique().tolist())
    sel_family = st.selectbox("Job Family", ["All"] + families)

    if sel_family != "All":
        df_view = df[df["Job Family"] == sel_family]
    else:
        df_view = df

    # Org chart
    html = build_network(df_view, None if sel_family == "All" else sel_family)
    components.html(html, height=640, scrolling=True)

    # Painel de detalhes
    st.markdown("### Detalhes do cargo")
    if sel_family != "All":
        sub_options = ["All"] + sorted(df_view["Sub Job Family"].dropna().unique().tolist())
        sel_sub = st.selectbox("Sub Job Family", sub_options)
        if sel_sub != "All":
            df_view = df_view[df_view["Sub Job Family"] == sel_sub]

    role_options = sorted(df_view["Job Profile"].dropna().unique().tolist())
    if not role_options:
        st.info("Nenhum cargo para mostrar.")
        return
    sel_role = st.selectbox("Job Profile", role_options)
    prof_row = df_view[df_view["Job Profile"] == sel_role].iloc[0].to_dict()
    attachments = match_ppt(ppts, prof_row["Job Family"], prof_row["Sub Job Family"], prof_row["Job Profile"])
    render_infocard(prof_row, attachments)


if __name__ == "__main__":
    main()
