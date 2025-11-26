import re
from pathlib import Path
from typing import Dict, List

import pandas as pd
import streamlit as st
from pptx import Presentation


@st.cache_data
def load_profiles(path: Path) -> pd.DataFrame:
    df = pd.read_excel(path).fillna("")
    needed = ["Job Family", "Sub Job Family", "Job Profile"]
    missing = [c for c in needed if c not in df.columns]
    if missing:
        st.error(f"Colunas ausentes no Job Profile.xlsx: {', '.join(missing)}")
        return pd.DataFrame(columns=needed)
    return df


@st.cache_data
def list_ppts() -> List[Path]:
    candidates = []
    for folder in ["PPTs", "ppts"]:
        d = Path(folder)
        if d.exists():
            candidates.extend(d.glob("*.pptx"))
    return sorted(candidates)


def match_ppt(ppts: List[Path], family: str, sub_family: str, profile: str) -> Path | None:
    tokens = [t.lower() for t in [profile, sub_family, family] if t]
    for ppt in ppts:
        name = ppt.stem.lower()
        if any(tok and tok in name for tok in tokens):
            return ppt
    return None


def extract_slides(path: Path) -> List[Dict[str, str]]:
    prs = Presentation(path)
    slides = []
    for slide in prs.slides:
        title, parts = "", []
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            txt = shape.text.strip()
            if not txt:
                continue
            if not title:
                title = txt
            else:
                parts.append(txt)
        if title or parts:
            slides.append({"title": title, "body": "\n\n".join(parts)})
    return slides


def render_canvas(slides: List[Dict[str, str]]) -> None:
    # canvas estilo cards lado a lado
    cols = st.columns(2)
    for idx, slide in enumerate(slides):
        col = cols[idx % 2]
        with col:
            col.markdown(
                f"""
                <div style="border:1px solid #e0ddd6; border-radius:12px; padding:14px; margin-bottom:12px; background:#ffffff; box-shadow:0 2px 6px rgba(0,0,0,0.06); min-height:140px;">
                    <div style="font-size:15px; font-weight:700; color:#0f4ad6; margin-bottom:6px;">{slide.get('title') or f'Slide {idx+1}'}</div>
                    <div style="font-size:13px; line-height:1.5; color:#222; white-space:pre-line;">{slide.get('body','—')}</div>
                </div>
                """,
                unsafe_allow_html=True,
            )


def main():
    st.set_page_config(page_title="Career Path", layout="wide")
    st.title("Career Path")
    st.markdown("Selecione a job family, subfamily e cargo para visualizar o infográfico (PPT convertido).")

    profiles = load_profiles(Path("data") / "Job Profile.xlsx")
    ppts = list_ppts()

    if profiles.empty:
        st.stop()

    families = sorted(profiles["Job Family"].dropna().unique().tolist())
    sel_fam = st.selectbox("Job Family", families)

    sub_options = profiles[profiles["Job Family"] == sel_fam]["Sub Job Family"].dropna().unique().tolist()
    sel_sub = st.selectbox("Sub Job Family", sub_options)

    prof_options = profiles[
        (profiles["Job Family"] == sel_fam) & (profiles["Sub Job Family"] == sel_sub)
    ]["Job Profile"].dropna().unique().tolist()
    sel_prof = st.selectbox("Job Profile", prof_options)

    ppt_path = match_ppt(ppts, sel_fam, sel_sub, sel_prof)
    if not ppt_path:
        st.warning("Nenhum PPT encontrado para este cargo (considerando os exemplos).")
        return

    slides = extract_slides(ppt_path)
    if not slides:
        st.info("Não foi possível extrair conteúdo deste PPT.")
        return

    st.markdown(f"### Infográfico · {sel_prof}")
    render_canvas(slides)


if __name__ == "__main__":
    main()
