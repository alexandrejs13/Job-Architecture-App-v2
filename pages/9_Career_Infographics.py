import re
from pathlib import Path
from typing import List, Dict

import streamlit as st
from pptx import Presentation


@st.cache_data
def list_ppts(ppts_dir: Path) -> List[Path]:
    return sorted(ppts_dir.glob("*.pptx"))


def infer_title_from_filename(path: Path) -> str:
    name = path.stem
    name = re.sub(r"^Career_Path[_-]?", "", name, flags=re.IGNORECASE)
    name = name.replace("_", " ").replace("-", " ")
    return name.strip().title()


def extract_slides(path: Path) -> List[Dict[str, str]]:
    prs = Presentation(path)
    slides = []
    for slide in prs.slides:
        parts = []
        title = ""
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if not text:
                continue
            if not title:
                title = text
            else:
                parts.append(text)
        if title or parts:
            slides.append({"title": title, "body": "\n\n".join(parts)})
    return slides


def render_infographic(slides: List[Dict[str, str]]) -> None:
    for idx, slide in enumerate(slides, start=1):
        st.markdown(
            f"""
            <div style="border:1px solid #e0ddd6; border-radius:12px; padding:16px; background:#ffffff; margin-bottom:14px; box-shadow:0 2px 6px rgba(0,0,0,0.05);">
                <div style="font-size:16px; font-weight:700; color:#0f4ad6; margin-bottom:8px;">Slide {idx}: {slide['title'] or 'Untitled'}</div>
                <div style="font-size:14px; line-height:1.5; white-space:pre-line; color:#222;">{slide['body'] or '—'}</div>
            </div>
            """,
            unsafe_allow_html=True,
        )


def main():
    st.set_page_config(page_title="Career Infographics", layout="wide")
    ppts_dir = Path("ppts")
    files = list_ppts(ppts_dir)

    st.title("Career Path Infographics")
    st.markdown("Selecione um cargo (PPT) para visualizar o conteúdo como infográfico.")

    if not files:
        st.warning("Nenhum PPT encontrado na pasta 'ppts'.")
        st.stop()

    options = {infer_title_from_filename(f): f for f in files}
    display_titles = sorted(options.keys())
    choice = st.selectbox("Selecione o cargo", display_titles)
    selected_path = options[choice]

    slides = extract_slides(selected_path)
    if not slides:
        st.info("Não foi possível extrair conteúdo desse PPT.")
        return

    st.markdown(f"### {choice}")
    render_infographic(slides)


if __name__ == "__main__":
    main()
